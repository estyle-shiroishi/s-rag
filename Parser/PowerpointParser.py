from Parser import IParser, chunk_text
from typing import List, Dict, Any
import logging
from pptx import Presentation
from io import BytesIO
from utils.aoai import aoai_chatgpt
import base64
import traceback

class PPTXChunkParser(IParser):
    """
    PPTXファイルをチャンクに変換するクラスです。
    パーサー用のインターフェイスであるIParserを継承しています。

    Attributes:
        __separator_charactor (str): チャンクを分割するためのセパレータ文字列です。
        __encoding (str): ファイルのエンコーディングです。
    """

    def __init__(self, separator="。", encoding="UTF-8"):
        """
        PPTXChunkParserのインスタンスを初期化します。

        Args:
            separator (str, optional): チャンクを分割するためのセパレータ文字列です。デフォルトは "。" です。
            encoding (str, optional): ファイルのエンコーディングです。デフォルトは "UTF-8" です。
        """
        self.__separator_charactor = separator
        self.__encoding = encoding
        self.SUPPORTED_MIME_TYPES = {
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg",
            "png": "image/png",
            "gif": "image/gif",
            "bmp": "image/bmp",
        }

    def parse(self, file_content: bytes) -> List[Dict[str, Any]]:
        """
        PPTXファイルをチャンクに変換するメソッドです。

        Args:
            file_content (bytes): ファイルの内容をバイト列で受け取ります。

        Returns:
            List[Dict[str, Any]]: ページ番号とテキストのリストを含む辞書のリストを返します。
        """
        logging.info("Parsing PPTX file.")

        # PPTXファイルの読み込み
        try: 
            prs = Presentation(BytesIO(file_content))

            # テキストと画像を抽出
            all_texts = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_texts.append(shape.text)
                    if shape.shape_type == 13:  # 画像の場合
                        try:
                            # 画像をAzure OpenAIに送信
                            image = shape.image
                            image_bytes = image.blob
                            image_base64 = self.__image_bytes_to_data_url(image_bytes, image.ext)
                            print(image.ext)
                            messages=[ # 4oに送るインプットメッセージ
                                            { "role": "system", "content": "You are a helpful assistant." }, #システムコンテント
                                            { "role": "user", "content": [ 
                                                {
                                                    "type": "text",
                                                    "text": "この画像の内容だけを出力してください" # GPT-4oに送るプロンプト（画像の内容を4oに出力してもらう）
                                                },
                                                {
                                                    "type": "image_url",
                                                    "image_url": {
                                                        "url": image_base64 # 画像のデータ
                                                    }
                                                }
                                            ] }
                                        ]
                            image_description = aoai_chatgpt(messages)
                            slide_texts.append(image_description)

                        except Exception as e:
                            # 画像の処理中にエラーが発生した場合はログに出力して処理を継続する
                            logging.error(f"Error occurred while processing image: {e}")
                            continue

                all_texts.append("".join(slide_texts))

            chunks = [chunk_text(text) for text in all_texts]

            # 空のチャンクを削除
            chunks_without_empty = [list(filter(None, chunk)) for chunk in chunks]

            page_with_chunk = [
                dict(page_number=i, texts=chunk)
                for i, chunk in enumerate(chunks_without_empty, start=1)
            ]
            
            if not page_with_chunk:
                logging.critical("parsing failed")
                raise ValueError(f"Failed to parse document.")
            
        except Exception as e:
            logging.error(f"An error occurred while parsing the Powerpoint file: {str(e)}")
            logging.debug(traceback.format_exc())
            return []
        
        return page_with_chunk

    def __image_bytes_to_data_url(self, image_bytes: bytes, ext: str) -> str:
        """
        画像のバイトデータと拡張子を受け取り、データURLを返します。

        Args:
            image_bytes: 画像のバイトデータ
            ext: 画像の拡張子 (例: "png", "jpg")

        Returns:
            データURL (例: "data:image/png;base64,...")

        Raises:
            ValueError: サポートされていない拡張子が指定された場合
        """
        mime_type = self.SUPPORTED_MIME_TYPES.get(ext.lower())
        if mime_type is None:
            raise ValueError(f"Unsupported image extension: {ext}")

        base64_encoded_data = base64.b64encode(image_bytes).decode("utf-8")
        return f"data:{mime_type};base64,{base64_encoded_data}"
